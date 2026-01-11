#!/usr/bin/env python3
"""
Create PowerPoint presentation from tutorial diagrams and narration scripts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# Colors from design system
DEEP_BLUE = RGBColor(0x09, 0x13, 0x58)
DARK_NAVY = RGBColor(0x0d, 0x08, 0x29)
STEEL_BLUE = RGBColor(0x68, 0x82, 0xa1)
BRIGHT_RED = RGBColor(0xb6, 0x00, 0x02)
WHITE = RGBColor(0xff, 0xff, 0xff)

# Slide dimensions for 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Image dimensions (A4 landscape ratio maintained)
IMG_WIDTH = Inches(11.5)
IMG_HEIGHT = Inches(8.14)  # Maintains A4 ratio

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide with dark navy background."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_NAVY
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(28)
        sub_para.font.color.rgb = STEEL_BLUE
        sub_para.alignment = PP_ALIGN.CENTER

    return slide

def create_section_slide(prs, part_num, part_title):
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DEEP_BLUE
    background.line.fill.background()

    # Part number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    num_frame = num_box.text_frame
    num_para = num_frame.paragraphs[0]
    num_para.text = f"PART {part_num}"
    num_para.font.size = Pt(36)
    num_para.font.color.rgb = STEEL_BLUE
    num_para.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = part_title
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    return slide

def create_diagram_slide(prs, image_path, speaker_notes=""):
    """Create a slide with a diagram image and speaker notes."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # White background (default)

    # Calculate centered position
    left = (SLIDE_WIDTH - IMG_WIDTH) / 2
    top = (SLIDE_HEIGHT - IMG_HEIGHT) / 2

    # Add image centered, scaled to fit
    if Path(image_path).exists():
        # Scale image to fit slide while maintaining aspect ratio
        pic = slide.shapes.add_picture(
            str(image_path),
            Inches(0.4),  # Left margin
            Inches(0.1),  # Top margin
            width=Inches(12.5),  # Fill width with margins
        )

    # Add speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = speaker_notes

    return slide

def create_end_slide(prs):
    """Create the final slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_NAVY
    background.line.fill.background()

    # Main message
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Now Go Build Something!"
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = "You have everything you need."
    sub_para.font.size = Pt(28)
    sub_para.font.color.rgb = STEEL_BLUE
    sub_para.alignment = PP_ALIGN.CENTER

    return slide

# Define presentation structure with diagrams and notes
PRESENTATION_STRUCTURE = [
    # Title
    {"type": "title", "title": "FROM NON-CODER TO BUILDER", "subtitle": "Coding with AI for Engineers"},

    # Part 1
    {"type": "section", "part": 1, "title": "FOUNDATIONS"},
    {"type": "diagram", "file": "01-1-the-core.png", "notes": """THE CORE PARTNERSHIP

You don't need to become a programmer. You need to become a builder.

Four elements form this partnership:

YOU — the engineer. You bring domain knowledge, judgment, and understanding of what you actually need.

AI — your coding partner. It knows syntax, patterns, and can generate code quickly.

PROMPT — the bridge. How you communicate with AI. Quality of prompt = quality of output.

SOFTWARE — the output. Working code that solves your problem.

Key insight: You provide direction. AI provides code. Together, you build.

You're not replacing yourself with AI. You're amplifying yourself."""},

    {"type": "diagram", "file": "02-1-the-loop.png", "notes": """THE LOOP

This is how every coding session with AI actually works. Memorize this loop.

DESCRIBE — Tell AI what you want.
GET — AI gives you code.
RUN — Actually run the code. See what happens.
EVALUATE — Does it work? Does it do what you wanted?

Then decide: Done or Refine?

If it works, move on. If not, go back to DESCRIBE with more information.

Key insight: Every problem is solvable through iteration. You don't need to get it right the first time."""},

    {"type": "diagram", "file": "02-2-the-draft.png", "notes": """THE DRAFT

Most beginners expect perfection on the first try. That's wrong.

Version 1: 10% — Rough draft, might not even run
Version 2: 40% — Basic structure works
Version 3: 75% — Right calculation, missing edge cases
Version 4: 100% — It works. Ship it.

Treat AI output like a first draft, not a final answer.

Stop asking "why didn't AI give me perfect code?" and start asking "what do I need to tell AI to make this better?" """},

    {"type": "diagram", "file": "03-1-six-operations.png", "notes": """SIX OPERATIONS

Everything you'll ever do with AI coding falls into six categories:

CREATE — Making something new
READ — Understanding existing code
EDIT — Changing existing code
RUN — Executing code (you do this)
FIX — Debugging errors
EXTEND — Adding to existing code

When stuck, ask yourself: "Which operation am I trying to do?"

Then structure your prompt around that operation."""},

    # Part 2
    {"type": "section", "part": 2, "title": "SKILLS"},
    {"type": "diagram", "file": "04-1-decomposition.png", "notes": """DECOMPOSITION

The single most important skill: breaking big things into small things.

Big task: "Build a bolt calculator app" — Too vague for AI

Break it down:
- Input section → Get force, Get area, Validate
- Calculation engine → Calculate stress, Check limits
- Results display → Format output, Show warnings

Big tasks confuse AI. Small tasks get perfect results.

Think of it like engineering drawings. You don't hand someone a napkin sketch. You provide detailed drawings of each component."""},

    {"type": "diagram", "file": "04-2-the-rule.png", "notes": """THE ONE-SENTENCE RULE

Simple test: Can I explain this task in ONE sentence?

If NO → Task is too big. Break it down.
If YES → Task is just right. Give it to AI.

TOO BIG: "Build me an app that handles all our bolt calculations with a UI and database and reports"

JUST RIGHT: "Create a function that calculates tensile stress from force and area"

When you find yourself writing a paragraph, stop. Decompose first."""},

    {"type": "diagram", "file": "05-1-the-spectrum.png", "notes": """THE SPECIFICITY SPECTRUM

Not all prompts are equal.

VAGUE: "Help me with calculations" — 20% success
MEDIUM: "Write a stress calculation function" — 50% success
PRECISE: "Write a Python function called calculate_tensile_stress that takes force in Newtons and area in mm², returns stress in MPa, raises error if negative" — 90% success

More specific prompts = Fewer iterations = Better results

Spend 2 minutes on a good prompt, get working code. Or spend 30 seconds on a vague prompt and 20 minutes fixing it."""},

    {"type": "diagram", "file": "05-2-five-components.png", "notes": """FIVE COMPONENTS OF A PRECISE PROMPT

WHAT — The action: "Create a function"
WHERE — The location: "in utils.py"
HOW — The method: "using recursion"
WHY — The purpose: "to calculate factorials for the stats module"
EXAMPLE — Show expected: "factorial(5) should return 120"

You don't always need all five. But the more you include, the better your results.

More components = Less guessing = Better code"""},

    {"type": "diagram", "file": "06-1-good-vs-bad.png", "notes": """GOOD VS BAD FEEDBACK

When AI gives you something that's not quite right, how you respond matters.

BAD FEEDBACK:
- "It doesn't work" — AI has no idea what's wrong
- "This is wrong" — Wrong how?
- "Fix it" — Fix what? AI will guess wrong
- "That's not what I wanted" — Then what DID you want?

GOOD FEEDBACK:
- "Getting error: TypeError on line 15" — AI knows where to look
- "Output is 5, but should be 10" — AI knows the gap
- "Change the loop to start at 1, not 0" — AI knows the fix

Specific feedback → Specific fixes → Done faster"""},

    {"type": "diagram", "file": "06-2-the-formula.png", "notes": """THE FEEDBACK FORMULA

"Got [X]. Expected [Y]. Change [Z]."

GOT [X] — What actually happened
EXPECTED [Y] — What should happen
CHANGE [Z] — What to modify (optional but helpful)

Example:
Got: "Function returns 15 for input [1,2,3,4,5]"
Expected: "Should return 120 (product, not sum)"
Change: "Use multiplication instead of addition in the loop"

Complete information = One-shot fix"""},

    {"type": "diagram", "file": "07-1-the-window.png", "notes": """THE WINDOW

AI only sees what's in the prompt.

Your world contains: project history, other files, business requirements, team standards, past conversations, your preferences.

AI's window contains: ONLY what you put in this prompt.

AI cannot see your screen, access your files, remember past chats, know your project, or read your mind. Unless you tell it.

The more context you share, the better AI understands."""},

    {"type": "diagram", "file": "07-2-what-to-share.png", "notes": """WHAT TO SHARE

Five types of context that improve AI results:

1. CONSTRAINTS — "Must work in Python 3.8", "No external libraries"
2. PREFERENCES — "Use descriptive variable names", "Prefer simple over clever"
3. HISTORY — "We tried X but it failed because...", "Part of a larger system"
4. ERRORS — Exact error messages, line numbers, when it happens
5. FILES — Current code, sample data, config files

Before sending: Did I share constraints? Preferences? History? Errors? Files?

More context = Fewer iterations = Better results"""},

    {"type": "diagram", "file": "08-1-the-decision.png", "notes": """THE DECISION

Sometimes you need to know when to stop refining and start fresh.

Is AI's output close to what you want?

YES → REFINE
Signs: Structure correct, small bugs, missing one feature, style needs adjustment

NO → RESTART
Signs: Wrong approach entirely, missing the point, multiple fundamental problems, going in circles

Starting fresh isn't failure. It's efficiency. Sometimes the fastest path is to abandon and try again with a better prompt."""},

    {"type": "diagram", "file": "08-2-the-5-attempt-rule.png", "notes": """THE 5-ATTEMPT RULE

Simple rule: Five attempts maximum, then restart.

Attempt 1: TRY — Initial prompt
Attempt 2: REFINE — Add specifics
Attempt 3: ADJUST — Fix errors
Attempt 4: REPHRASE — Different approach
Attempt 5: LAST TRY — Everything you've learned

Still not there? RESTART with a new conversation and better prompt.

Too few (1-2): Giving up too early
Just right (3-5): Enough to find solutions
Too many (6+): Diminishing returns, going in circles"""},

    # Part 3
    {"type": "section", "part": 3, "title": "PROBLEMS"},
    {"type": "diagram", "file": "09-1-common-traps.png", "notes": """COMMON TRAPS

TRAP 1: BLIND TRUST — Running code without reading it
Fix: Always review before running

TRAP 2: VAGUE PROMPTS — "Make it better"
Fix: Be specific about what's wrong

TRAP 3: TOO BIG TASKS — "Build me a whole app"
Fix: Break into small tasks

TRAP 4: NO CONTEXT — Not sharing errors or code
Fix: Paste errors, include code

TRAP 5: GIVING UP TOO FAST — First attempt didn't work
Fix: Give 3-5 attempts

TRAP 6: NOT LEARNING — Copy-paste without understanding
Fix: Ask AI to explain

Most traps come from treating AI like magic instead of a tool."""},

    {"type": "diagram", "file": "09-2-trust-levels.png", "notes": """TRUST LEVELS

How much to trust AI output depends on consequences of failure.

LOW TRUST — Verify Everything
When: Security code, financial calculations, safety-critical, data deletion
Action: Read every line, test thoroughly, get second opinion

MEDIUM TRUST — Review and Test
When: Standard features, business logic, API integrations
Action: Skim code, run basic tests, check edge cases

HIGH TRUST — Quick Check
When: Simple utilities, formatting, text manipulation
Action: Glance at output, run once, move on

Trust level depends on consequences, not AI confidence.
Higher stakes = More verification needed"""},

    {"type": "diagram", "file": "10-1-the-debugging.png", "notes": """THE DEBUGGING FLOW

1. CAPTURE — Copy the exact error message
2. LOCATE — Find the problem line
3. CONTEXT — Include surrounding code
4. ASK — Send to AI with all context
5. APPLY — Make the fix
6. TEST — Verify it's fixed

BAD: "My code doesn't work, fix it"
GOOD: "Error: TypeError on line 15. Here's lines 10-20. Input was [1, 2, None]. Should skip None values."

Better information in = Better fix out"""},

    {"type": "diagram", "file": "10-2-error-types.png", "notes": """ERROR TYPES

SYNTAX ERROR — Code can't even run
Examples: Missing brackets, typos, wrong indentation
Tell AI: "Syntax error on line X"
Easiest to fix

RUNTIME ERROR — Code crashes while running
Examples: Division by zero, null access, file not found
Tell AI: "Crashes at line X with [error]"
Medium difficulty

LOGIC ERROR — Code runs but wrong result
Examples: Wrong calculation, off-by-one, wrong condition
Tell AI: "Got X, expected Y"
Hardest to fix

Name the error type = Faster fix"""},

    {"type": "diagram", "file": "11-1-the-stuck.png", "notes": """GETTING STUCK

Warning signs:
- Same error keeps coming back
- Fixing one thing breaks another
- AI gives same answer repeatedly
- More than 5 attempts without progress
- Getting frustrated

What's happening:
- Wrong mental model of the problem
- Missing crucial information
- Task is too complex
- Solving wrong problem entirely

The stuck loop: TRY → FAIL → TRY AGAIN → SAME FAIL → repeat

Being stuck is information: your approach needs to change."""},

    {"type": "diagram", "file": "11-2-escape-routes.png", "notes": """ESCAPE ROUTES

1. START FRESH — New conversation, new prompt
   Use when: Same error repeating

2. BREAK SMALLER — Split the task
   Use when: Task seems impossible

3. ADD CONTEXT — Share more information
   Use when: AI seems confused

4. REPHRASE — Different words, different angle
   Use when: AI misunderstanding you

5. ASK WHY — Get AI to explain the error
   Use when: You don't understand

6. SIMPLIFY — Make a minimal example
   Use when: Too many things could be wrong

Every problem has an escape route."""},

    # Part 4
    {"type": "section", "part": 4, "title": "BUILDING"},
    {"type": "diagram", "file": "12-1-learning-to-read.png", "notes": """LEARNING TO READ CODE

You don't need to write code to understand it.

LOOK FOR:
- Names: Variables describe what they hold
- Structure: Indentation shows grouping
- Keywords: if, for, while, return
- Comments: Human explanations
- Flow: Top to bottom, branching

ASK AI:
- "What does this code do?"
- "Explain line by line"
- "What is [variable] for?"
- "Why is this needed?"

SKIP FOR NOW:
- Memorizing syntax
- Understanding every detail
- Advanced patterns

Goal: Understand what code does, not memorize how to write it."""},

    {"type": "diagram", "file": "12-2-code-anatomy.png", "notes": """CODE ANATOMY

Universal parts of every program:

COMMENT — Notes for humans, ignored by computer
Like sticky notes on code

FUNCTION — Reusable blocks of code
Take inputs, give outputs
Like machines in a factory

VARIABLE — Named containers for data
The name tells you what's inside

CONDITION — Makes decisions
if/else, true/false
Like a fork in the road

OUTPUT — Shows results to user
Makes the invisible visible

Every program is made of these building blocks."""},

    {"type": "diagram", "file": "13-1-variables.png", "notes": """VARIABLES

Variables are named containers for data. Like labeled boxes.

bolt_diameter = 12.5
Creates a box labeled "bolt_diameter" with 12.5 inside.

Types:
- NUMBERS: count = 42, price = 19.99
- TEXT: name = "M12", status = "OK"
- TRUE/FALSE: is_valid = True
- LISTS: sizes = [8, 10, 12]

Naming rules:
- No spaces (use underscores)
- Start with letter
- Be descriptive

Good: bolt_count, max_stress
Bad: x, bc, thing1

Good variable names make code self-documenting."""},

    {"type": "diagram", "file": "13-2-control-flow.png", "notes": """CONTROL FLOW

Three patterns handle most logic:

IF/ELSE — Decisions
if stress > 250:
    print("Fail")
else:
    print("Pass")

FOR LOOP — Repeat N times
for size in [8, 10, 12]:
    print(size)

WHILE LOOP — Repeat until done
while error > 0.01:
    refine()

Summary:
- if/else: Choose between options
- for: Do something N times
- while: Keep going until done

Programs = Data + Decisions + Repetition"""},

    {"type": "diagram", "file": "14-1-functions.png", "notes": """FUNCTIONS

A function is like a machine: put something in, get something out.

def calculate_stress(force, area):
    return force / area

result = calculate_stress(1000, 4)
# result is now 250

Parts:
- Name: what it does
- Parameters: inputs
- Body: the work
- Return: output

Why use functions:
- Reuse code (write once)
- Organize logic (clear names)
- Test easily (isolated units)
- Fix in one place

Functions = Reusable building blocks"""},

    {"type": "diagram", "file": "14-2-libraries.png", "notes": """LIBRARIES

Libraries are pre-built code. Don't reinvent the wheel.

WITHOUT: 5 lines to calculate average manually
WITH: 1 line using statistics.mean()

Useful libraries:

Math/Science: numpy, scipy, pandas, math
Visualization: matplotlib, plotly, seaborn
Files: os, json, csv, datetime

Ask AI: "What library should I use to [task]?"

Stand on shoulders of giants — use libraries!"""},

    {"type": "diagram", "file": "15-1-data-structures.png", "notes": """DATA STRUCTURES

LIST — Ordered collection
bolt_sizes = [8, 10, 12]
Access by position: bolt_sizes[0]
Like a shopping list

DICTIONARY — Key-value pairs
bolt = {"name": "M12", "grade": "8.8"}
Access by name: bolt["grade"]
Like a real dictionary

TUPLE — Fixed, unchangeable
coords = (100.0, 200.5)
Can't change after creation
Like a sealed envelope

Use LIST for collections that change
Use DICT for named properties
Use TUPLE for fixed values

Right structure = Cleaner code"""},

    {"type": "diagram", "file": "15-2-files-io.png", "notes": """FILES AND I/O

Files are permanent storage. Data survives when program ends.

READING:
with open("data.csv") as file:
    content = file.read()

WRITING:
with open("results.txt", "w") as file:
    file.write("Stress: 250 MPa")

File types:
- .txt, .csv, .json — Human readable
- .xlsx, .xml, .db — Use libraries

Ask AI: "Read this CSV into a list" or "Write results to file"

Files = Permanent data that survives restarts"""},

    {"type": "diagram", "file": "16-1-testing.png", "notes": """TESTING

Testing is quality control for code.

Types:
- Manual: Run, check output, try different inputs
- Automated: Write test code, runs automatically
- Edge cases: Empty inputs, extreme values

Simple test:
result = calculate_stress(1000, 4)
expected = 250
assert result == expected

PASS = Code works
FAIL = Fix needed

Ask AI: "Write tests for this function"

Tested code = Trusted code"""},

    {"type": "diagram", "file": "16-2-project-structure.png", "notes": """PROJECT STRUCTURE

Typical structure:
bolt_calculator/
    main.py           # Entry point
    calculations.py   # Core functions
    utils.py          # Helpers
    data/             # Input files
    output/           # Results
    tests/            # Test files
    README.md         # Documentation

Why organize:
- Find files quickly
- Separate concerns
- Easy to share
- Professional standard

File naming: lowercase, underscores, descriptive

Ask AI: "Create project structure for [description]"

Organized project = Maintainable project"""},

    # Part 5
    {"type": "section", "part": 5, "title": "APPLICATION"},
    {"type": "diagram", "file": "17-1-your-first-project.png", "notes": """YOUR FIRST PROJECT

Start small, make it work, then expand.

Phase 1: PLAN
Define goal, inputs/outputs, simplest version
Write it in one sentence before coding

Phase 2: BUILD
Start with core feature, test as you go, add features one by one
Get it working before making it perfect

Phase 3: REFINE
Handle edge cases, add error messages, clean up
Ask AI to review

Example: Bolt Stress Calculator
Plan: "Calculate stress from force and area, warn if too high"
Build: Create function, add input, add warning, format output
Refine: Handle zero area, add units, save to file

Done is better than perfect — ship V1, improve later"""},

    {"type": "diagram", "file": "17-2-project-checklist.png", "notes": """PROJECT CHECKLIST

Before Starting:
□ Goal clearly defined?
□ Inputs and outputs identified?
□ Broken into small tasks?

While Building:
□ Testing each piece?
□ Understanding AI's code?
□ Saving working versions?

Before Finishing:
□ Works with normal input?
□ Works with edge cases?
□ Has clear error messages?

Quality Check:
□ Code readable?
□ Names descriptive?
□ Comments where needed?

Final Questions:
- Can someone else use this without help?
- Would you trust this in production?

Checklist complete = Ready to ship"""},

    {"type": "diagram", "file": "18-1-growing-skills.png", "notes": """GROWING YOUR SKILLS

Stage 1: COPY
Follow AI's code exactly, ask for explanations, learn patterns
Goal: Get code working, understand what it does

Stage 2: MODIFY
Adjust AI's code yourself, combine pieces, fix small issues
Goal: Adapt solutions, build on existing code

Stage 3: CREATE
Design solutions, use AI as assistant, guide the process
Goal: Build anything, AI accelerates but doesn't replace you

The secret: Every expert was once a beginner who kept building.

Growth = Practice + Curiosity + Persistence"""},

    {"type": "diagram", "file": "18-2-resources.png", "notes": """RESOURCES

AI Tools:
- Claude — Great for explanations
- ChatGPT — Good for code generation
- GitHub Copilot — In-editor assistance

Learning More:
- Python.org — Documentation
- Stack Overflow — Q&A for errors
- YouTube — Visual learning

Practice Ideas:
- Automate calculations
- Data visualization
- File processing

When Stuck:
1. Ask AI → 2. Search Web → 3. Rephrase → 4. Break Down

Most problems solved in 1-2 steps.

Help is always one question away."""},

    {"type": "diagram", "file": "18-3-the-journey.png", "notes": """THE JOURNEY

You started: "I can't code"

You learned:
- The Loop: Describe → Get → Run → Evaluate
- Decomposition: Big → Small tasks
- Communication: Specific prompts → Better results

You are now: "I build with AI"

Skills you have:
- Read and understand code
- Communicate with AI effectively
- Debug and fix problems
- Build working solutions

What's next:
- Pick a real problem
- Build your first project
- Share it with someone
- Build another one!

The truth: Everyone looks things up. Everyone gets stuck. Everyone asks for help.
The only difference is persistence."""},

    {"type": "diagram", "file": "18-4-summary.png", "notes": """KEY TAKEAWAYS

YOU + AI = BUILDER

Six Principles:
1. THE LOOP — Iterate until it works
2. DECOMPOSE — One sentence = right size
3. BE SPECIFIC — What, where, how, why, example
4. GIVE FEEDBACK — Got X, Expected Y, Change Z
5. SHARE CONTEXT — AI only knows what you tell it
6. KNOW WHEN TO RESTART — 5 attempts then fresh start

The Formula:
Clear Goal + Small Steps + Iteration + Persistence = Working Code

You now have everything you need.
Start building."""},

    # End
    {"type": "end"},
]

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    diagrams_path = Path("/home/veritarium/dev/projects/tutorial/diagrams/source")

    for item in PRESENTATION_STRUCTURE:
        if item["type"] == "title":
            create_title_slide(prs, item["title"], item.get("subtitle", ""))
        elif item["type"] == "section":
            create_section_slide(prs, item["part"], item["title"])
        elif item["type"] == "diagram":
            image_path = diagrams_path / item["file"]
            create_diagram_slide(prs, image_path, item.get("notes", ""))
        elif item["type"] == "end":
            create_end_slide(prs)

    # Save presentation
    output_path = Path("/home/veritarium/dev/projects/tutorial/From-NonCoder-to-Builder.pptx")
    prs.save(str(output_path))
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
